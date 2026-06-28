"""
PPTX presentation parser using python-pptx.
Extracts text content from presentations.
"""

from pptx import Presentation
import logging
from typing import Dict, List, Any

logger = logging.getLogger(__name__)


class PPTXParser:
    """
    Parser for PPTX presentations.
    Extracts text content from slides.
    """

    @staticmethod
    def parse(file_path: str, max_length: int = 50000) -> Dict[str, Any]:
        """
        Parse a PPTX file and extract its content.
        
        Args:
            file_path: Path to the PPTX file
            max_length: Maximum characters to extract (default: 50000)
        
        Returns:
            Dictionary containing:
                - text: Extracted text content
                - slides: Number of slides
                - metadata: Presentation metadata
                - success: Boolean indicating parsing success
                - error: Error message if parsing failed
        """
        try:
            prs = Presentation(file_path)

            # Extract metadata
            metadata = {
                'title': prs.core_properties.title or 'Unknown',
                'author': prs.core_properties.author or 'Unknown',
                'subject': prs.core_properties.subject or 'Unknown',
                'slides': len(prs.slides)
            }

            # Extract text from all slides
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_content.append(f"--- Slide {slide_num} ---\n" + '\n'.join(slide_text))

            # Combine all extracted text
            combined_text = '\n\n'.join(text_content)

            # Truncate if necessary
            if len(combined_text) > max_length:
                combined_text = combined_text[:max_length] + '\n\n[Content truncated...]'

            return {
                'success': True,
                'text': combined_text,
                'tables': [],  # PPTX tables would require more complex parsing
                'metadata': metadata,
                'page_count': metadata['slides'],
                'has_tables': False,
                'error': None
            }

        except FileNotFoundError:
            error_msg = f"PPTX file not found: {file_path}"
            logger.error(error_msg)
            return {
                'success': False,
                'text': '',
                'tables': [],
                'metadata': {},
                'page_count': 0,
                'has_tables': False,
                'error': error_msg
            }

        except Exception as e:
            error_msg = f"Error parsing PPTX: {str(e)}"
            logger.error(error_msg)
            return {
                'success': False,
                'text': '',
                'tables': [],
                'metadata': {},
                'page_count': 0,
                'has_tables': False,
                'error': error_msg
            }
